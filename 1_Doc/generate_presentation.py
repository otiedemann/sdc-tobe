#!/usr/bin/env python3
"""Generate SDC-ToBe Tech Stack PPTX presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ────────────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD    = RGBColor(0x22, 0x22, 0x3A)
ACCENT     = RGBColor(0x00, 0xB4, 0xD8)
ACCENT2    = RGBColor(0x48, 0xCA, 0xE4)
GREEN      = RGBColor(0x06, 0xD6, 0xA0)
ORANGE     = RGBColor(0xFF, 0xB7, 0x03)
RED        = RGBColor(0xEF, 0x47, 0x6F)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
MID_GRAY   = RGBColor(0x88, 0x88, 0x99)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Smaller corner radius
    shape.adjustments[0] = 0.05
    return shape


def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(4)

        # Bullet char
        run_b = p.add_run()
        run_b.text = "\u25B8 "
        run_b.font.size = Pt(font_size)
        run_b.font.color.rgb = bullet_color
        run_b.font.name = "Calibri"

        run_t = p.add_run()
        run_t.text = item
        run_t.font.size = Pt(font_size)
        run_t.font.color.rgb = color
        run_t.font.name = "Calibri"
    return txBox


def add_table_slide(slide, left, top, width, rows_data, col_widths=None,
                    header_color=ACCENT, font_size=14):
    """Add a styled table. rows_data[0] = headers."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.4 * n_rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Calibri"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = LIGHT_GRAY
            # Cell fill
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = RGBColor(0x00, 0x76, 0x8A)
            elif r % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(0x28, 0x28, 0x42)
            else:
                cell.fill.fore_color.rgb = RGBColor(0x20, 0x20, 0x38)
    return table_shape


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide, BG_DARK)

# Accent line
add_rect(slide, Inches(1), Inches(2.6), Inches(2), Pt(4), ACCENT)

add_text(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.2),
         "SDC-ToBe", font_size=54, color=WHITE, bold=True)
add_text(slide, Inches(1), Inches(2.0), Inches(11), Inches(0.8),
         "Swarm Drone Challenge 2026", font_size=28, color=ACCENT)
add_text(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
         "Technical Architecture & Stack Overview", font_size=22, color=LIGHT_GRAY)
add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.6),
         "Multi-drone swarm control platform with real-time ArUco positioning,\n"
         "autonomous navigation, and competitive game strategy",
         font_size=16, color=MID_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — System Overview / Architecture
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
         "System Architecture", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.0), Inches(1.5), Pt(3), ACCENT)

# Layer boxes
layers = [
    ("C2 Strategy Dashboard", "Game AI, phase management, scoring\nFastAPI + WebSocket (port 9090)", ORANGE),
    ("C2 Interface", "Web UI, grid map, mission control\nFastAPI + WebSocket (port 8080)", ACCENT),
    ("Unified API Server", "REST API, RC control, video, telemetry\nFlask on Raspberry Pi (port 8080)", GREEN),
    ("Drone Hardware / Simulator", "DJI Tello | Parrot Anafi | Sim\nWiFi 2.4GHz, Olympe/djitellopy SDK", RED),
]

y = Inches(1.5)
for title, desc, color in layers:
    card = add_rect(slide, Inches(0.8), y, Inches(7.5), Inches(1.15), BG_CARD, color)
    add_text(slide, Inches(1.1), y + Pt(8), Inches(3.5), Inches(0.5),
             title, font_size=18, color=color, bold=True)
    add_text(slide, Inches(4.5), y + Pt(8), Inches(3.6), Inches(1.0),
             desc, font_size=13, color=LIGHT_GRAY)
    y += Inches(1.35)

# Side box: ArUco Positioning
add_rect(slide, Inches(8.8), Inches(1.5), Inches(3.8), Inches(5.4), BG_CARD, ACCENT2)
add_text(slide, Inches(9.1), Inches(1.6), Inches(3.3), Inches(0.5),
         "ArUco Positioning", font_size=18, color=ACCENT2, bold=True)
add_bullet_list(slide, Inches(9.0), Inches(2.2), Inches(3.5), Inches(4.5), [
    "25 wall-mounted markers",
    "OpenCV ArUco detection",
    "solvePnP pose estimation",
    "Kalman filter smoothing",
    "20m x 10m arena grid",
    "Markers at 2m + 4m height",
    "18x18 cm marker size",
    "SSE + UDP position relay",
], font_size=13, color=LIGHT_GRAY, bullet_color=ACCENT2)

# Arrow annotations
add_text(slide, Inches(5.0), Inches(6.4), Inches(4), Inches(0.5),
         "WiFi 2.4GHz  |  HTTP REST  |  WebSocket  |  SSE  |  UDP",
         font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Tech Stack
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
         "Tech Stack", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.0), Inches(1.5), Pt(3), ACCENT)

# Cards
cards = [
    ("Backend", ACCENT, [
        "Python 3.11+",
        "FastAPI + Uvicorn (C2)",
        "Flask (Pi controllers)",
        "Pydantic data models",
        "Threading / async-await",
    ]),
    ("Vision & Sensing", GREEN, [
        "OpenCV 4.13 (ArUco)",
        "Camera calibration (NPZ)",
        "solvePnP 3D localization",
        "Kalman filter (1D/3D)",
        "Power-law distance calib.",
    ]),
    ("Drone SDKs", ORANGE, [
        "djitellopy 2.5 (Tello)",
        "Parrot Olympe (Anafi)",
        "Auto-detect by IP addr",
        "ARM-safe conditional import",
        "Gimbal + GPS (Anafi)",
    ]),
    ("Frontend & Comms", RED, [
        "Vanilla JS + CSS Grid",
        "WebSocket (real-time)",
        "SSE telemetry streams",
        "MJPEG video streaming",
        "UDP position relay",
    ]),
]

x = Inches(0.5)
for title, color, items in cards:
    add_rect(slide, x, Inches(1.5), Inches(2.9), Inches(4.5), BG_CARD, color)
    add_text(slide, x + Inches(0.2), Inches(1.6), Inches(2.5), Inches(0.5),
             title, font_size=20, color=color, bold=True)
    add_bullet_list(slide, x + Inches(0.15), Inches(2.2), Inches(2.6), Inches(3.5),
                    items, font_size=14, color=LIGHT_GRAY, bullet_color=color)
    x += Inches(3.1)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Supported Drones
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
         "Supported Drones", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.0), Inches(1.5), Pt(3), ACCENT)

drones = [
    ("DJI Tello", ACCENT, [
        "SDK: djitellopy 2.5",
        "Control: UDP socket",
        "Default IP: 192.168.10.1",
        "Video: 960x720 MJPEG",
        "Telemetry: battery, height, IMU",
        "Multi-drone: per-WiFi-stick isolation",
    ]),
    ("Parrot Anafi", GREEN, [
        "SDK: Parrot Olympe (x86 only)",
        "Control: Olympe PCMD + moveBy",
        "Default IP: 192.168.42.1",
        "Video: 1280x720 (up to 4K)",
        "Sensors: GPS, gimbal, barometer",
        "Features: RTH, geofence, flight config",
    ]),
    ("Simulator", ORANGE, [
        "4 variants (2D, 3D, API, C2)",
        "Drop-in API replacement",
        "Physics: velocity model",
        "3D web viewer (Three.js)",
        "Configurable arena & targets",
        "No hardware required",
    ]),
]

x = Inches(0.5)
for title, color, items in drones:
    add_rect(slide, x, Inches(1.5), Inches(3.9), Inches(5.0), BG_CARD, color)
    add_text(slide, x + Inches(0.3), Inches(1.65), Inches(3.3), Inches(0.5),
             title, font_size=22, color=color, bold=True)
    add_bullet_list(slide, x + Inches(0.2), Inches(2.3), Inches(3.5), Inches(4.0),
                    items, font_size=15, color=LIGHT_GRAY, bullet_color=color)
    x += Inches(4.1)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Unified API Server
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
         "Unified API Server", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.0), Inches(1.5), Pt(3), ACCENT)

add_text(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6),
         "Single REST endpoint auto-detects drone type (Tello vs Anafi) by IP address",
         font_size=16, color=LIGHT_GRAY)

# API table
rows = [
    ["Endpoint", "Method", "Description"],
    ["/api/takeoff", "POST", "Take off"],
    ["/api/land", "POST", "Land"],
    ["/api/rc", "POST", "Raw RC control (lr, fb, ud, yaw: -100..100)"],
    ["/api/move", "POST", "Discrete move (direction + cm)"],
    ["/api/telemetry", "GET", "Full telemetry snapshot (JSON)"],
    ["/api/telemetry/stream", "GET", "SSE telemetry stream (~2.5 Hz)"],
    ["/api/video/start", "POST", "Start MJPEG video stream"],
    ["/api/video", "GET", "MJPEG stream output"],
    ["/api/position", "GET", "Current ArUco position"],
    ["/api/position/events", "GET", "SSE position updates"],
    ["/api/position/video", "GET", "MJPEG with ArUco overlay"],
    ["/api/settings", "GET/POST", "Flight limits (persisted to JSON)"],
    ["/api/heartbeat", "GET", "Watchdog keepalive"],
]
add_table_slide(slide, Inches(0.8), Inches(1.9), Inches(11.5), rows,
                col_widths=[Inches(3.2), Inches(1.3), Inches(7.0)], font_size=12)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ArUco Positioning System
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
         "ArUco Positioning System", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.0), Inches(1.5), Pt(3), ACCENT)

# Pipeline
steps = [
    ("1. Detection", "OpenCV detects\nArUco 4x4 markers\nin each video frame"),
    ("2. Pose Estimation", "solvePnP computes\n3D camera position\nfrom marker corners"),
    ("3. Kalman Filter", "Smooths noisy\nmeasurements for\nstable positioning"),
    ("4. Relay", "SSE + UDP streams\nposition to C2\nat ~2.5 Hz"),
]

x = Inches(0.5)
for title, desc in steps:
    add_rect(slide, x, Inches(1.5), Inches(2.8), Inches(2.2), BG_CARD, ACCENT)
    add_text(slide, x + Inches(0.2), Inches(1.6), Inches(2.4), Inches(0.5),
             title, font_size=17, color=ACCENT, bold=True)
    add_text(slide, x + Inches(0.2), Inches(2.2), Inches(2.4), Inches(1.3),
             desc, font_size=14, color=LIGHT_GRAY)
    x += Inches(3.05)

# Arena details
add_rect(slide, Inches(0.5), Inches(4.0), Inches(5.8), Inches(3.0), BG_CARD, GREEN)
add_text(slide, Inches(0.8), Inches(4.1), Inches(5.4), Inches(0.5),
         "Arena Configuration", font_size=18, color=GREEN, bold=True)
add_bullet_list(slide, Inches(0.7), Inches(4.7), Inches(5.4), Inches(2.2), [
    "20m wide x 10m deep arena",
    "25 markers (ID 0-24) on 4 walls",
    "Front, Back, Left, Right walls",
    "Each marker at 2m or 4m height",
    "Physical marker size: 18x18 cm",
    "ArUco dictionary: 4x4_100",
], font_size=14, color=LIGHT_GRAY, bullet_color=GREEN)

# Calibration
add_rect(slide, Inches(6.6), Inches(4.0), Inches(5.8), Inches(3.0), BG_CARD, ORANGE)
add_text(slide, Inches(6.9), Inches(4.1), Inches(5.4), Inches(0.5),
         "Distance Calibration", font_size=18, color=ORANGE, bold=True)
add_bullet_list(slide, Inches(6.8), Inches(4.7), Inches(5.4), Inches(2.2), [
    "Power-law: dist = 109.17 * px^(-0.897)",
    "Empirical fit from measured data",
    "55px = 3.0m, 85px = 2.0m, 180px = 1.0m",
    "Camera intrinsics via calibration.npz",
    "Client-side ArUco (VideoMarkerTracker)",
    "Perspective skew for angle detection",
], font_size=14, color=LIGHT_GRAY, bullet_color=ORANGE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Autonomous Tools: ArUco Seek
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
         "Autonomous Tool: ArUco Seek & Approach", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.0), Inches(1.5), Pt(3), ACCENT)

# State machine phases
phases = [
    ("TAKEOFF", "Climb to\ntarget altitude"),
    ("SETTLE", "Stabilize\n(3 seconds)"),
    ("SCAN", "Rotate 360\xb0\nfind markers"),
    ("ALIGN", "Yaw to face\nmarker dead-on"),
    ("APPROACH", "Fly forward\nto target dist"),
    ("HOVER", "Hold position\n1.5m from marker"),
]

x = Inches(0.3)
colors = [ACCENT, ACCENT, GREEN, ORANGE, ORANGE, RED]
for i, (name, desc) in enumerate(phases):
    add_rect(slide, x, Inches(1.5), Inches(1.9), Inches(1.6), BG_CARD, colors[i])
    add_text(slide, x + Inches(0.1), Inches(1.6), Inches(1.7), Inches(0.4),
             name, font_size=15, color=colors[i], bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), Inches(2.1), Inches(1.7), Inches(0.9),
             desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    x += Inches(2.05)

# Features
add_rect(slide, Inches(0.5), Inches(3.5), Inches(5.8), Inches(3.5), BG_CARD, ACCENT)
add_text(slide, Inches(0.8), Inches(3.6), Inches(5.4), Inches(0.5),
         "Visual Servoing", font_size=18, color=ACCENT, bold=True)
add_bullet_list(slide, Inches(0.7), Inches(4.2), Inches(5.4), Inches(2.6), [
    "Pure image-based control (no world coordinates)",
    "Marker center in image -> yaw + altitude correction",
    "Pixel size -> distance estimate (calibration curve)",
    "Perspective skew -> lateral strafe correction",
    "Edge-boost: 3x gain when marker nears frame edge",
    "P-controllers: yaw P=40, alt P=45, dist P=8",
    "Works with ANY room setup (no arena config needed)",
], font_size=14, color=LIGHT_GRAY, bullet_color=ACCENT)

add_rect(slide, Inches(6.6), Inches(3.5), Inches(5.8), Inches(3.5), BG_CARD, GREEN)
add_text(slide, Inches(6.9), Inches(3.6), Inches(5.4), Inches(0.5),
         "Safety & CLI", font_size=18, color=GREEN, bold=True)
add_bullet_list(slide, Inches(6.8), Inches(4.2), Inches(5.4), Inches(2.6), [
    "Ctrl+C: abort and land immediately",
    "Watchdog heartbeat every 400ms",
    "Auto-land on 120s timeout",
    "Auto-land if marker lost for 8s",
    "--marker <ID>: target specific marker",
    "--hover-distance: set standoff distance",
    "--record: save debug video (MJPEG)",
], font_size=14, color=LIGHT_GRAY, bullet_color=GREEN)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — C2 Interfaces
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
         "Command & Control Interfaces", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.0), Inches(1.5), Pt(3), ACCENT)

# C2 Interface
add_rect(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.3), BG_CARD, ACCENT)
add_text(slide, Inches(0.8), Inches(1.6), Inches(5.4), Inches(0.5),
         "C2 Interface (Team Control)", font_size=20, color=ACCENT, bold=True)
add_bullet_list(slide, Inches(0.7), Inches(2.3), Inches(5.4), Inches(4.3), [
    "FastAPI + WebSocket (port 8080)",
    "20m x 10m arena grid (200 cells)",
    "Zone system: Red | Neutral | Blue",
    "Click-to-navigate: select drone, click cell",
    "Per-drone MJPEG video panels",
    "Mission flow: navigate -> search -> hover -> return",
    "Modes: live (real drones), simulator, sim_api",
    "Multi-WiFi isolation for 2-5 real Tellos",
    "Target designation (up to 6 per team)",
], font_size=14, color=LIGHT_GRAY, bullet_color=ACCENT)

# C2 Strategy
add_rect(slide, Inches(6.6), Inches(1.5), Inches(5.8), Inches(5.3), BG_CARD, ORANGE)
add_text(slide, Inches(6.9), Inches(1.6), Inches(5.4), Inches(0.5),
         "C2 Strategy (Game AI)", font_size=20, color=ORANGE, bold=True)
add_bullet_list(slide, Inches(6.8), Inches(2.3), Inches(5.4), Inches(4.3), [
    "FastAPI + WebSocket (port 9090)",
    "Rule-based strategy engine",
    "Game phases:",
    "  SETUP -> SCOUTING -> ATTACK",
    "  -> TEAM_RETURN -> DEFEND",
    "  -> INSTANT_WIN -> GAME_OVER",
    "Drone roles: Scout, Attacker,",
    "  Defender, Returner",
    "Live ArUco position via UDP relay",
    "Scoring logic + win detection",
], font_size=14, color=LIGHT_GRAY, bullet_color=ORANGE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Simulation Environments
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
         "Simulation Environments", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.0), Inches(1.5), Pt(3), ACCENT)

sims = [
    ("2D Swarm Sim", "sim_swarm/", ACCENT, [
        "Lightweight (no ext. deps)",
        "Scenario-based testing",
        "baseline, degraded_link, dropout",
        "Metrics: time, collisions, battery",
        "Timeline replay in web viewer",
    ]),
    ("3D API Simulator", "sim_swarm_API/", GREEN, [
        "Full HTTP + WebSocket API",
        "Three.js 3D web visualization",
        "Drop-in for unified API",
        "Swarm mission logic",
        "ArUco marker simulation",
    ]),
    ("HTTP/UDP Drone Sim", "drone-sim/", ORANGE, [
        "Minimal physics model",
        "Tello-compatible HTTP API",
        "UDP JSON telemetry output",
        "Single-file (14KB)",
        "Configurable targets",
    ]),
    ("C2-Integrated Sim", "sim_swarm_c2/", RED, [
        "Full game simulation",
        "Strategy engine overlay",
        "Multi-drone coordination",
        "Phase + role testing",
        "Scoring validation",
    ]),
]

x = Inches(0.3)
for title, path, color, items in sims:
    add_rect(slide, x, Inches(1.5), Inches(3.0), Inches(5.0), BG_CARD, color)
    add_text(slide, x + Inches(0.2), Inches(1.6), Inches(2.6), Inches(0.5),
             title, font_size=17, color=color, bold=True)
    add_text(slide, x + Inches(0.2), Inches(2.1), Inches(2.6), Inches(0.3),
             path, font_size=11, color=MID_GRAY)
    add_bullet_list(slide, x + Inches(0.15), Inches(2.5), Inches(2.7), Inches(3.8),
                    items, font_size=13, color=LIGHT_GRAY, bullet_color=color)
    x += Inches(3.15)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Communication & Data Flow
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
         "Communication & Data Flow", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.0), Inches(1.5), Pt(3), ACCENT)

rows = [
    ["Channel", "Protocol", "Frequency", "Direction", "Purpose"],
    ["Drone Control", "HTTP REST", "On-demand", "C2 -> Pi -> Drone", "Takeoff, land, move, RC"],
    ["Telemetry", "SSE (HTTP)", "~2.5 Hz", "Drone -> Pi -> C2", "Battery, height, yaw, status"],
    ["Video Feed", "MJPEG (HTTP)", "15-30 FPS", "Drone -> Pi -> Client", "Live camera stream"],
    ["ArUco Position", "SSE + UDP", "~2.5 Hz", "Pi -> C2 / Tools", "x, y, z, velocity, markers"],
    ["C2 Dashboard", "WebSocket", "Real-time", "Bidirectional", "State updates, commands"],
    ["RC Commands", "HTTP POST", "5 Hz", "Tool -> Pi", "lr, fb, ud, yaw (-100..100)"],
    ["Heartbeat", "HTTP GET", "2.5 Hz", "Tool -> Pi", "Watchdog keepalive"],
    ["Video Record", "Server-side", "Frame rate", "Pi local", "MP4 with ArUco overlay"],
]
add_table_slide(slide, Inches(0.5), Inches(1.5), Inches(12.3), rows,
                col_widths=[Inches(2.0), Inches(1.6), Inches(1.5), Inches(2.8), Inches(4.4)],
                font_size=13)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Key Numbers
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
         "By the Numbers", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.0), Inches(1.5), Pt(3), ACCENT)

stats = [
    ("~67", "Python\nModules", ACCENT),
    ("10", "Major\nSystems", GREEN),
    ("3", "Drone\nTypes", ORANGE),
    ("4", "Simulator\nVariants", RED),
    ("25", "ArUco\nMarkers", ACCENT2),
    ("40+", "REST API\nEndpoints", GREEN),
    ("200", "Arena Grid\nCells (1m\u00b2)", ORANGE),
    ("18cm", "Marker\nSize", RED),
]

x = Inches(0.4)
y = Inches(1.8)
for i, (num, label, color) in enumerate(stats):
    if i == 4:
        x = Inches(0.4)
        y = Inches(4.2)
    card = add_rect(slide, x, y, Inches(2.8), Inches(2.0), BG_CARD, color)
    add_text(slide, x, y + Inches(0.2), Inches(2.8), Inches(0.9),
             num, font_size=44, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(1.1), Inches(2.8), Inches(0.7),
             label, font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    x += Inches(3.05)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Project Structure
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
         "Project Structure", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.0), Inches(1.5), Pt(3), ACCENT)

rows = [
    ["Directory", "Purpose", "Key Tech"],
    ["controller_unified/", "Unified Pi API server (Tello + Anafi)", "Flask, OpenCV, Olympe"],
    ["c2_interface/", "Web-based team C2 (grid map, missions)", "FastAPI, WebSocket, JS"],
    ["c2_strategy/", "Strategic game AI + scoring", "FastAPI, state machine"],
    ["aruco-position/", "ArUco detection + 3D localization", "OpenCV, solvePnP, Kalman"],
    ["tools/", "Autonomous drone tools (aruco_seek)", "Visual servoing, P-control"],
    ["sim_swarm/", "Lightweight 2D swarm simulation", "Pure Python, scenarios"],
    ["sim_swarm_API/", "Full 3D simulator + web viewer", "FastAPI, Three.js"],
    ["drone-sim/", "Minimal HTTP/UDP drone simulator", "Flask, physics model"],
    ["sim_swarm_c2/", "C2-integrated game simulation", "Strategy + sim combined"],
    ["1_Doc/", "Documentation, specs, calibration", "PDF, DOCX, XLSX"],
]
add_table_slide(slide, Inches(0.5), Inches(1.3), Inches(12.3), rows,
                col_widths=[Inches(2.5), Inches(5.5), Inches(4.3)], font_size=13)


# ═══════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════
out_path = "/Users/tiedemann/GitHub/sdc-tobe/1_Doc/SDC-ToBe_Tech_Stack.pptx"
prs.save(out_path)
print(f"Saved to {out_path}")
